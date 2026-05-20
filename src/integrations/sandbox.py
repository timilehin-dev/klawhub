import time
import hmac
import hashlib
import json
import logging
import httpx
from typing import Dict, Any, Optional
from src.config import settings

logger = logging.getLogger("klawhub.sandbox")

class SandboxClient:
    def __init__(self, function_url: Optional[str] = None, webhook_secret: Optional[str] = None):
        self.function_url = function_url or settings.modal_function_url
        self.webhook_secret = (webhook_secret or settings.modal_webhook_secret).encode('utf-8')

    def _generate_headers(self, payload_str: str) -> Dict[str, str]:
        """Generates HMAC-SHA256 signature and timestamp headers to secure sandbox execution."""
        timestamp = str(int(time.time()))
        
        # Sign payload + timestamp using HMAC SHA-256
        message = f"{payload_str}:{timestamp}".encode('utf-8')
        signature = hmac.new(self.webhook_secret, message, hashlib.sha256).hexdigest()
        
        return {
            "X-Webhook-Timestamp": timestamp,
            "X-Webhook-Signature": signature,
            "X-Webhook-Secret": self.webhook_secret.decode('utf-8'),
            "Content-Type": "application/json"
        }

    async def execute_code(
        self, 
        code: str, 
        language: str = "python", 
        timeout_seconds: int = 120,
        dependencies: Optional[list[str]] = None,
        memory_tier: Optional[str] = None,
        workspace_id: Optional[Any] = None
    ) -> Dict[str, Any]:
        """Executes a dynamic code script inside the secure, isolated Modal sandbox.
        
        Intelligently auto-detects dependencies and memory requirements based on import analysis.
        Automatically checks and mounts custom skills cached in the database for the given workspace.
        Returns a dict containing stdout, stderr, execution duration, and exit status code.
        """
        deps = list(dependencies) if dependencies else []
        tier = memory_tier or "standard"
        mounted_skills = {}

        if language == "python":
            import ast
            import sys
            
            try:
                tree = ast.parse(code)
                imported = set()
                for node in ast.walk(tree):
                    if isinstance(node, ast.Import):
                        for name in node.names:
                            imported.add(name.name.split('.')[0])
                    elif isinstance(node, ast.ImportFrom):
                        if node.module:
                            imported.add(node.module.split('.')[0])

                # standard libraries to ignore
                stdlib = getattr(sys, "stdlib_module_names", set())

                # pre-installed packages in Modal
                pre_installed = {
                    "fastapi", "requests", "httpx", "lxml", "beautifulsoup4", "polars", "pandas",
                    "numpy", "matplotlib", "seaborn", "plotly", "scikit-learn", "typst", "pypandoc",
                    "markdown", "pdfplumber", "weasyprint", "fastembed", "crawl4ai", "lightpanda-py",
                    "playwright", "json", "math", "datetime", "time", "re", "csv", "collections", "itertools"
                }

                # heavy modules that trigger the 16GB tier
                heavy_modules = {
                    "torch", "tensorflow", "transformers", "scipy", "fastembed", "spacy", "crawl4ai",
                    "weasyprint", "playwright", "polars"
                }

                # --- Auto-Detect and Mount Custom Skills ---
                if workspace_id:
                    try:
                        import uuid
                        from sqlmodel import select
                        from src.db.pool import get_db_session
                        from src.db.models import Skill

                        normalized_workspace_id = uuid.UUID(str(workspace_id))
                        async with get_db_session() as session:
                            stmt = select(Skill).where(Skill.workspace_id == normalized_workspace_id, Skill.is_active == True)
                            db_skills = (await session.execute(stmt)).scalars().all()
                            skills_map = {s.name.lower().strip(): s for s in db_skills}
                            
                            # --- Dynamic First-Time Skill Installation Blocker ---
                            imported_lowered = [mod.lower().strip() for mod in imported]
                            if "document_generator" in imported_lowered and "document_generator" not in skills_map:
                                logger.info("document_generator skill requested for the first time in this workspace. Installing from GitHub...")
                                from src.core.evolution.acquisition import SkillAcquisitionEngine
                                
                                repo_url = "https://github.com/timilehin-dev/klawhub-standard-skills"
                                success = await SkillAcquisitionEngine.clone_and_register_github_skill(
                                    workspace_id=normalized_workspace_id,
                                    repo_url=repo_url,
                                    file_path="document_generator.py",
                                    skill_name="document_generator",
                                    description="Industry standard professional document generator (PDF, DOCX, XLSX, PPTX)"
                                )
                                
                                # Highly resilient fallback: if GitHub fetch fails, register the standard implementation directly
                                if not success:
                                    logger.warning("GitHub fetch failed or repository not found. Activating failsafe high-fidelity local cache...")
                                    
                                    fallback_code = """# Dynamic Document Generator Skill for Klawhub
# Provides high-fidelity generation of PDF, DOCX, XLSX, and PPTX documents inside the isolated sandbox.
import os
import sys

def generate_pdf(html_content: str, filename: str) -> str:
    \"\"\"Generates a professional PDF from HTML using WeasyPrint.\"\"\"
    import weasyprint
    weasyprint.HTML(string=html_content).write_pdf(filename)
    return os.path.abspath(filename)

def generate_docx(markdown_content: str, filename: str) -> str:
    \"\"\"Generates a professional DOCX from Markdown using pypandoc/python-docx.\"\"\"
    import pypandoc
    try:
        pypandoc.convert_text(markdown_content, 'docx', format='md', outputfile=filename)
    except Exception as e:
        import docx
        doc = docx.Document()
        for line in markdown_content.splitlines():
            if line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("* ") or line.startswith("- "):
                doc.add_paragraph(line[2:], style='List Bullet')
            else:
                doc.add_paragraph(line)
        doc.save(filename)
    return os.path.abspath(filename)

def generate_xlsx(data_matrix: list, filename: str, headers: list = None) -> str:
    \"\"\"Generates a professional Excel spreadsheet with styled tables using openpyxl.\"\"\"
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Exported Data"
    
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
    center_align = Alignment(horizontal="center", vertical="center")
    border_side = Side(border_style="thin", color="D9D9D9")
    thin_border = Border(left=border_side, right=border_side, top=border_side, bottom=border_side)
    
    current_row = 1
    if headers:
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=current_row, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = center_align
            cell.border = thin_border
        ws.row_dimensions[current_row].height = 25
        current_row += 1
        
    for row_data in data_matrix:
        for col_idx, val in enumerate(row_data, 1):
            cell = ws.cell(row=current_row, column=col_idx, value=val)
            cell.border = thin_border
            if isinstance(val, (int, float)):
                cell.number_format = '#,##0.00' if isinstance(val, float) else '#,##0'
        ws.row_dimensions[current_row].height = 18
        current_row += 1
        
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = openpyxl.utils.get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
        
    wb.save(filename)
    return os.path.abspath(filename)

def generate_pptx(slides_data: list, filename: str) -> str:
    \"\"\"Generates a PowerPoint presentation using python-pptx.\"\"\"
    import pptx
    prs = pptx.Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    for slide_info in slides_data:
        title = slide_info.get("title", "")
        content = slide_info.get("content", [])
        is_title_slide = slide_info.get("is_title", False)
        
        if is_title_slide:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title
            if content:
                slide.placeholders[1].text = "\\n".join(content) if isinstance(content, list) else content
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.clear()
            if isinstance(content, list):
                for idx, bullet in enumerate(content):
                    if idx == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
            else:
                tf.paragraphs[0].text = content
                
    prs.save(filename)
    return os.path.abspath(filename)
"""
                                    new_skill = Skill(
                                        workspace_id=normalized_workspace_id,
                                        name="document_generator",
                                        description="Industry standard professional document generator (PDF, DOCX, XLSX, PPTX)",
                                        category="custom",
                                        repo_url=repo_url,
                                        file_path="document_generator.py",
                                        entrypoint="generate_pdf",
                                        source_code=fallback_code,
                                        dependencies="weasyprint,pypandoc,openpyxl,python-docx,python-pptx",
                                        is_active=True
                                    )
                                    session.add(new_skill)
                                    await session.commit()
                                    logger.info("Successfully registered failsafe document_generator skill inside database.")
                                    
                                # Reload database skills map
                                stmt = select(Skill).where(Skill.workspace_id == normalized_workspace_id, Skill.is_active == True)
                                db_skills = (await session.execute(stmt)).scalars().all()
                                skills_map = {s.name.lower().strip(): s for s in db_skills}
                            
                            for mod in list(imported):
                                normalized_mod = mod.lower().strip()
                                if normalized_mod in skills_map:
                                    skill = skills_map[normalized_mod]
                                    logger.info(f"Auto-detect mounted skill requirement: '{normalized_mod}'")
                                    
                                    # Add to mounted skills dictionary
                                    mounted_skills[normalized_mod] = {
                                        "code": skill.source_code,
                                        "dependencies": [d.strip() for d in skill.dependencies.split(",")] if skill.dependencies else []
                                    }
                                    
                                    # If the skill itself has dynamic dependencies, append them
                                    if skill.dependencies:
                                        for dep in skill.dependencies.split(","):
                                            dep = dep.strip()
                                            if dep and dep not in deps:
                                                deps.append(dep)
                                    
                                    # If the skill or its dependencies import heavy packages, promote memory tier
                                    if skill.dependencies:
                                        for dep in skill.dependencies.split(","):
                                            dep_base = dep.split("==")[0].strip().lower()
                                            if dep_base in heavy_modules:
                                                tier = "heavy"
                                                
                                    # Exclude the skill itself from pip package installation list!
                                    if mod in imported:
                                        imported.remove(mod)
                    except Exception as skill_err:
                        logger.warning(f"Failed to auto-detect/load workspace skills from database: {skill_err}")

                for mod in imported:
                    # Upgrade memory tier if a heavy module is imported
                    if mod in heavy_modules:
                        tier = "heavy"
                    
                    # Auto-detect dynamic packages to install (not standard, not pre-installed, and not local app modules)
                    if (mod not in stdlib and 
                        mod not in pre_installed and 
                        mod not in {"src", "api", "db", "workflows"} and 
                        mod not in deps):
                        # Simple rule to avoid appending standard built-in mock modules
                        if not mod.startswith("_"):
                            deps.append(mod)

            except Exception as e:
                logger.warning(f"Failed to auto-detect dependencies or memory tier: {e}")

        payload = {
            "code": code,
            "language": language,
            "timeout": timeout_seconds,
            "dependencies": deps,
            "memory_tier": tier,
            "mounted_skills": mounted_skills
        }
        payload_str = json.dumps(payload)
        headers = self._generate_headers(payload_str)
        
        logger.info(f"Dispatching dynamic sandbox job to Modal (language: {language}, timeout: {timeout_seconds}s)")
        
        async with httpx.AsyncClient(timeout=float(timeout_seconds + 10)) as client:
            try:
                response = await client.post(
                    self.function_url,
                    content=payload_str,
                    headers=headers
                )
                
                if response.status_code != 200:
                    logger.error(f"Modal sandbox returned error HTTP status: {response.status_code}")
                    return {
                        "success": False,
                        "exit_code": -1,
                        "stdout": "",
                        "stderr": f"Sandbox execution HTTP error: {response.status_code}\nContent: {response.text}",
                        "duration_ms": 0
                    }
                
                result = response.json()
                success = result.get("success") if "success" in result else (result.get("exit_code", -1) == 0)
                logger.info(f"Sandbox run completed successfully: {success}")
                return {
                    "success": success,
                    "exit_code": result.get("exit_code", 0 if success else -1),
                    "stdout": result.get("stdout", ""),
                    "stderr": result.get("stderr", ""),
                    "duration_ms": result.get("duration_ms", 0),
                    "error": result.get("error", None),
                    "generated_files": result.get("generated_files", [])
                }
                
            except httpx.RequestError as e:
                logger.exception("Failed to connect to Modal sandbox gateway")
                return {
                    "success": False,
                    "exit_code": -1,
                    "stdout": "",
                    "stderr": f"Network request to Sandbox failed: {str(e)}",
                    "duration_ms": 0
                }

# Global Sandbox Client instance
sandbox_client = SandboxClient()
